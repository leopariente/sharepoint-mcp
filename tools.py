import io
import re
import sys
from pathlib import Path
from mcp.server.fastmcp import FastMCP
from cache import _with_cache, get_stats as _cache_get_stats
from constants import (
    SHAREPOINT_SITE_URL,
    SHAREPOINT_BASE_FOLDER,
    STUDENT_EMAIL,
    STUDENT_PASSWORD,
    _REGISTRY,
)
from decorators import _with_backpressure
from office365.runtime.auth.user_credential import UserCredential
from office365.sharepoint.client_context import ClientContext
from office365.sharepoint.search.service import SearchService

_CLAUDE_MD = (Path(__file__).parent / "claude.md").read_text(encoding="utf-8")

mcp = FastMCP("sharepoint-agent", host="0.0.0.0", instructions=_CLAUDE_MD)


# ---------------------------------------------------------------------------
# SharePoint context — created once and reused across all tool calls.
# ---------------------------------------------------------------------------
_ctx: ClientContext | None = None


def _get_ctx() -> ClientContext:
    global _ctx
    if _ctx is not None:
        return _ctx
    if not STUDENT_PASSWORD:
        raise RuntimeError(
            "SHAREPOINT_PASSWORD environment variable is not set. "
            "Set it before starting the server."
        )
    credentials = UserCredential(STUDENT_EMAIL, STUDENT_PASSWORD)
    _ctx = ClientContext(SHAREPOINT_SITE_URL).with_credentials(credentials)
    _ctx.load(_ctx.web)
    _ctx.execute_query()
    return _ctx


# ---------------------------------------------------------------------------
# Tool: registry_lookup
# ---------------------------------------------------------------------------
@mcp.tool()
def registry_lookup(query: str) -> str:
    """Look up a course in the local registry to get its SharePoint path and year/semester.

    ALWAYS call this tool first when the user mentions a course by name.
    Only fall back to list_files / search_files if this returns no match.

    Args:
        query: Course name or partial name to search for (case-insensitive).

    Returns:
        Matching course(s) with name, year, semester, and full SharePoint path,
        or a "not found" message if no match exists.
    """
    q = query.lower().strip()
    matches = [
        c for c in _REGISTRY
        if q in c["name"].lower() or c["name"].lower() in q
    ]
    if not matches:
        return (
            f"No course matching '{query}' found in the registry. "
            "Use list_files or search_files to browse SharePoint."
        )
    lines = [f"Registry matches for '{query}':\n"]
    for c in matches:
        sem = f"Semester {c['semester']}" if c["semester"] else "—"
        year_label = f"Year {c['year']}" if c["year"] else "—"
        full_path = f"{SHAREPOINT_BASE_FOLDER}/{c['path']}"
        lines.append(f"  Course:   {c['name']}")
        lines.append(f"  Year:     {year_label}  |  Semester: {sem}")
        lines.append(f"  Path:     {full_path}")
        lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Tool: list_files
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def list_files(folder_url: str = SHAREPOINT_BASE_FOLDER) -> str:
    """List files and subdirectories in a SharePoint folder.

    Args:
        folder_url: Server-relative URL of the folder to list
                    (e.g. "/sites/MySite/Shared Documents/subfolder").
                    Leave empty to list the root of the default document library.

    Returns:
        A formatted listing of files (with sizes) and subfolders.
    """
    try:
        ctx = _get_ctx()
        folder = ctx.web.get_folder_by_server_relative_url(folder_url)
        ctx.load(folder, ["Files", "Folders"])
        ctx.execute_query()

        lines: list[str] = [f"Contents of {folder_url}:\n"]

        for f in folder.files:
            lines.append(f"  [FILE] {f.name}  ({f.length} bytes)")

        for sub in folder.folders:
            if sub.name == "Forms":
                continue
            lines.append(f"  [DIR]  {sub.name}/  -> {sub.serverRelativeUrl}")

        if len(lines) == 1:
            lines.append("  (empty)")

        return "\n".join(lines)
    except Exception as e:
        return _format_error("list_files", folder_url, e)


# ---------------------------------------------------------------------------
# Tool: read_file_content
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def read_file_content(file_url: str) -> str:
    """Read and return the text content of a file stored in SharePoint.

    Suitable for text-based files (.txt, .md, .csv, .json, .xml, .html, .py, etc.).
    Binary files will return an error message instead.

    Args:
        file_url: Server-relative URL of the file
                  (e.g. "/sites/MySite/Shared Documents/notes.txt").

    Returns:
        The text content of the file, or an error message.
    """
    try:
        ctx = _get_ctx()
        response = ctx.web.get_file_by_server_relative_url(file_url)
        buf = io.BytesIO()
        response.download(buf).execute_query()
        buf.seek(0)
        raw = buf.read()

        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return raw.decode("latin-1")
            except Exception:
                return (
                    f"[binary content — {len(raw)} bytes] "
                    "The file does not appear to be text-based."
                )
    except Exception as e:
        return _format_error("read_file_content", file_url, e)


# ---------------------------------------------------------------------------
# Tool: search_files
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def search_files(query: str, folder_url: str = "") -> str:
    """Search for files by name within a SharePoint folder (recursive).

    Args:
        query: Substring to match against file names (case-insensitive).
        folder_url: Server-relative URL of the folder to search in.
                    Leave empty to search the default document library.

    Returns:
        A list of matching file paths, or a message if none were found.
    """
    if not folder_url:
        folder_url = SHAREPOINT_BASE_FOLDER

    try:
        ctx = _get_ctx()
        search = SearchService(ctx)
        kql = f'filename:"{query}" path:"{SHAREPOINT_SITE_URL}"'
        result = search.post_query(
            query_text=kql,
            select_properties=["Path", "FileName"],
            row_limit=500,
            trim_duplicates=False,
        )
        ctx.execute_query()

        matches: list[str] = []
        for row in result.value.PrimaryQueryResult.RelevantResults.Table.Rows:
            cells = {c.Key: c.Value for c in row.Cells}
            path = cells.get("Path", "")
            if path:
                matches.append(path)

        if not matches:
            return f"No files matching '{query}' found under {folder_url}."

        header = f"Files matching '{query}' under {folder_url}:\n"
        return header + "\n".join(f"  {m}" for m in matches)
    except Exception as e:
        return _format_error("search_files", folder_url, e)


# ---------------------------------------------------------------------------
# Tool: search_content
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def search_content(
    query: str,
    folder_url: str = "",
    file_types: list[str] = ["pdf"],
    max_results: int = 10,
) -> str:
    """Full-text search inside SharePoint file contents.

    Args:
        query: Text to search for inside file contents (e.g. "dynamic programming recurrence").
        folder_url: Narrow the search to a specific folder (optional).
        file_types: File extensions to include (e.g. ["pdf", "docx"]). Default: ["pdf"].
        max_results: Maximum number of results to return. Default: 10.

    Returns:
        Ranked list of matching files with relevance score and a content snippet.
    """
    try:
        ctx = _get_ctx()
        search = SearchService(ctx)

        kql = f'"{query}"'
        if folder_url:
            host = SHAREPOINT_SITE_URL.split("/sites/")[0]
            kql += f' path:"{host}{folder_url}"'
        else:
            kql += f' path:"{SHAREPOINT_SITE_URL}"'
        if file_types:
            type_filter = " OR ".join(f"FileExtension:{t}" for t in file_types)
            kql += f" ({type_filter})"

        result = search.post_query(
            query_text=kql,
            select_properties=["Title", "Path"],
            row_limit=max_results,
            trim_duplicates=False,
        )
        ctx.execute_query()

        rows = result.value.PrimaryQueryResult.RelevantResults.Table.Rows
        if not rows:
            return f"No files found containing '{query}'."

        def _cells_to_dict(row):
            cells = {}
            for c in row.Cells:
                if hasattr(c, "Key") and hasattr(c, "Value"):
                    cells[c.Key] = c.Value
                elif isinstance(c, dict):
                    cells[c.get("Key", "")] = c.get("Value", "")
            return cells

        lines = [f"Files containing '{query}':\n"]
        for i, row in enumerate(rows, 1):
            cells = _cells_to_dict(row)
            title = cells.get("Title") or cells.get("FileName", "Unknown")
            path = cells.get("Path", "")

            lines.append(f"[{i}] {title}")
            lines.append(f"    Path:  {path}")
            lines.append("")

        return "\n".join(lines)
    except Exception as e:
        return _format_error("search_content", query, e)


# ---------------------------------------------------------------------------
# Shared helper
# ---------------------------------------------------------------------------
def _detect_language(text: str) -> str:
    has_hebrew = any("\u0590" <= c <= "\u05FF" for c in text)
    has_latin = any("a" <= c.lower() <= "z" for c in text)
    if has_hebrew and has_latin:
        return "Hebrew + English"
    return "Hebrew" if has_hebrew else "English"


# ---------------------------------------------------------------------------
# Tool: read_pdf
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def read_pdf(
    file_url: str,
    pages: str = "all",
    max_chars: int = 8000,
) -> str:
    """Extract text from a PDF stored in SharePoint.

    Handles Hebrew RTL text and mathematical notation properly.

    Args:
        file_url: SharePoint server-relative URL of the PDF file.
        pages: Page range to read — "all", a single page ("3"), or a range ("1-5").
        max_chars: Maximum characters to return. Default: 8000.

    Returns:
        Extracted text with page headers, file info, and detected language.
    """
    def _parse_page_range(pages_str: str, total: int) -> list[int]:
        s = pages_str.strip().lower()
        if s == "all":
            return list(range(total))
        if "-" in s:
            start, end = s.split("-", 1)
            return list(range(int(start) - 1, min(int(end), total)))
        return [int(s) - 1]

    try:
        import pymupdf
    except ImportError:
        return "[read_pdf] pymupdf is not installed. Run: pip install pymupdf"

    try:
        ctx = _get_ctx()
        file_obj = ctx.web.get_file_by_server_relative_url(file_url)
        buf = io.BytesIO()
        file_obj.download(buf).execute_query()
        buf.seek(0)
        raw = buf.read()

        doc = pymupdf.open(stream=io.BytesIO(raw), filetype="pdf")
        page_indices = _parse_page_range(pages, len(doc))

        extracted = []
        for i in page_indices:
            if 0 <= i < len(doc):
                text = doc[i].get_text("text", flags=pymupdf.TEXT_PRESERVE_LIGATURES)
                extracted.append(f"--- Page {i + 1} ---\n{text}")

        full_text = "\n".join(extracted)
        lang = _detect_language(full_text)
        size_mb = round(len(raw) / 1024 / 1024, 2)
        filename = file_url.split("/")[-1]

        header = (
            f"File: {filename}\n"
            f"Pages: {len(doc)} | Language: {lang} | Size: {size_mb} MB\n\n"
        )
        return header + full_text[:max_chars]
    except Exception as e:
        return _format_error("read_pdf", file_url, e)


# ---------------------------------------------------------------------------
# Tool: read_docx
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def read_docx(
    file_url: str,
    max_chars: int = 8000,
) -> str:
    """Extract text from a Word document (.docx) stored in SharePoint.

    Extracts paragraphs and table contents in reading order.

    Args:
        file_url: SharePoint server-relative URL of the .docx file.
        max_chars: Maximum characters to return. Default: 8000.

    Returns:
        Extracted text with heading markers and table contents, or an error message.
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return "[read_docx] python-docx is not installed. Run: pip install python-docx"

    try:
        ctx = _get_ctx()
        file_obj = ctx.web.get_file_by_server_relative_url(file_url)
        buf = io.BytesIO()
        file_obj.download(buf).execute_query()
        buf.seek(0)

        doc = Document(buf)
        filename = file_url.split("/")[-1]
        lines: list[str] = []

        for block in doc.element.body:
            tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

            if tag == "p":
                from docx.text.paragraph import Paragraph
                para = Paragraph(block, doc)
                text = para.text.strip()
                if not text:
                    continue
                style = para.style.name if para.style else ""
                if style.startswith("Heading"):
                    level = style.replace("Heading", "").strip()
                    prefix = "#" * int(level) if level.isdigit() else "##"
                    lines.append(f"{prefix} {text}")
                else:
                    lines.append(text)

            elif tag == "tbl":
                from docx.table import Table
                table = Table(block, doc)
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append(" | ".join(cells))
                lines.append("")

        full_text = "\n".join(lines)
        lang = _detect_language(full_text)
        header = (
            f"File: {filename}\n"
            f"Paragraphs: {len(doc.paragraphs)} | Language: {lang}\n\n"
        )
        return header + full_text[:max_chars]
    except Exception as e:
        return _format_error("read_docx", file_url, e)


# ---------------------------------------------------------------------------
# Tool: read_pptx
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def read_pptx(
    file_url: str,
    slides: str = "all",
    max_chars: int = 8000,
) -> str:
    """Extract text from a PowerPoint file (.pptx) stored in SharePoint.

    Extracts all text shapes from each slide, including titles and bullet points.

    Args:
        file_url: SharePoint server-relative URL of the .pptx file.
        slides: Slides to read — "all", a single slide ("3"), or a range ("1-5").
        max_chars: Maximum characters to return. Default: 8000.

    Returns:
        Extracted slide text with slide headers, or an error message.
    """
    def _parse_slide_range(slides_str: str, total: int) -> list[int]:
        s = slides_str.strip().lower()
        if s == "all":
            return list(range(total))
        if "-" in s:
            start, end = s.split("-", 1)
            return list(range(int(start) - 1, min(int(end), total)))
        return [int(s) - 1]

    try:
        from pptx import Presentation
    except ImportError:
        return "[read_pptx] python-pptx is not installed. Run: pip install python-pptx"

    try:
        ctx = _get_ctx()
        file_obj = ctx.web.get_file_by_server_relative_url(file_url)
        buf = io.BytesIO()
        file_obj.download(buf).execute_query()
        buf.seek(0)

        prs = Presentation(buf)
        filename = file_url.split("/")[-1]
        slide_indices = _parse_slide_range(slides, len(prs.slides))

        lines: list[str] = []
        for i in slide_indices:
            if 0 <= i < len(prs.slides):
                slide = prs.slides[i]
                lines.append(f"--- Slide {i + 1} ---")
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                lines.append("")

        full_text = "\n".join(lines)
        lang = _detect_language(full_text)
        header = (
            f"File: {filename}\n"
            f"Slides: {len(prs.slides)} | Language: {lang}\n\n"
        )
        return header + full_text[:max_chars]
    except Exception as e:
        return _format_error("read_pptx", file_url, e)


# ---------------------------------------------------------------------------
# Tool: get_file_metadata
# ---------------------------------------------------------------------------
@mcp.tool()
@_with_cache
@_with_backpressure
def get_file_metadata(
    file_url: str,
    include_versions: bool = False,
) -> str:
    """Return rich metadata for a SharePoint file.

    Args:
        file_url: SharePoint server-relative URL of the file.
        include_versions: Whether to include version history. Default: False.

    Returns:
        Author, dates, size, grade parsed from filename, and optional version history.
    """
    def _extract_grade(filename: str) -> str | None:
        match = re.search(r"_(\d{2,3})\.\w+$", filename, re.IGNORECASE)
        return match.group(1) if match else None

    try:
        ctx = _get_ctx()
        file_obj = ctx.web.get_file_by_server_relative_url(file_url)
        ctx.load(file_obj, ["Name", "Length", "TimeCreated", "TimeLastModified", "UIVersionLabel"])
        ctx.execute_query()

        author = "Unknown"
        modified_by = "Unknown"
        try:
            list_item = file_obj.listItemAllFields
            ctx.load(list_item, ["Author", "Editor"])
            ctx.execute_query()

            def _extract_name(val):
                if isinstance(val, dict):
                    return val.get("Title", val.get("LoginName", "Unknown"))
                return str(val) if val else "Unknown"

            author_val = list_item.get_property("Author")
            editor_val = list_item.get_property("Editor")
            modified_by = _extract_name(editor_val)
            author = _extract_name(author_val)
            if author == "Unknown":
                author = modified_by
        except Exception:
            pass

        props = file_obj.properties
        name = props.get("Name", file_url.split("/")[-1])
        size_mb = round(int(props.get("Length", 0)) / 1024 / 1024, 2)
        created = props.get("TimeCreated", "Unknown")
        modified = props.get("TimeLastModified", "Unknown")
        version = props.get("UIVersionLabel", "Unknown")
        grade = _extract_grade(name)

        lines = [
            f"File: {name}",
            "=" * 40,
            f"Path:     {file_url}",
            f"Author:      {author}",
            f"Modified By: {modified_by}",
            f"Uploaded:    {created}",
            f"Modified:    {modified}",
            f"Size:     {size_mb} MB",
            f"Grade:    {grade if grade else 'N/A (not in filename)'}",
            f"Version:  {version}",
        ]

        if include_versions:
            try:
                versions = file_obj.versions
                ctx.load(versions)
                ctx.execute_query()
                lines.append("\nVersion History:")
                for v in versions:
                    vp = v.properties
                    lines.append(f"  v{vp.get('VersionLabel', '?')} — {vp.get('Created', '?')}")
            except Exception as ve:
                lines.append(f"\nVersion history unavailable: {ve}")

        return "\n".join(lines)
    except Exception as e:
        return _format_error("get_file_metadata", file_url, e)


# ---------------------------------------------------------------------------
# Tool: cache_stats
# ---------------------------------------------------------------------------
@mcp.tool()
def cache_stats() -> str:
    """Return Redis cache hit/miss/bypass counters for this server process.

    Returns:
        JSON-formatted stats: hits, misses, bypasses, bytes_served, redis_healthy.
    """
    import json as _json
    return _json.dumps(_cache_get_stats(), indent=2)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _format_error(tool: str, target: str, exc: Exception) -> str:
    msg = str(exc)
    if "403" in msg or "Forbidden" in msg:
        return (
            f"[{tool}] 403 Forbidden for '{target}'. "
            "The current user does not have permission to access this resource."
        )
    if "404" in msg or "Not Found" in msg:
        return (
            f"[{tool}] 404 Not Found for '{target}'. "
            "Check that the path exists and is spelled correctly."
        )
    return f"[{tool}] Error accessing '{target}': {msg}"


# ---------------------------------------------------------------------------
# Eager auth check — called by both entry points on startup
# ---------------------------------------------------------------------------
def warmup() -> None:
    try:
        _get_ctx()
        print("SharePoint context ready.", file=sys.stderr)
    except Exception as e:
        print(f"Warning: eager auth failed ({e}), will retry on first tool call.", file=sys.stderr)
